import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_demo_ppt():
    prs = Presentation()
    
    # 16:9 Aspect Ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # SliceMatic Brand Colors
    BRAND_GREEN = RGBColor(34, 197, 94)
    DARK_BG = RGBColor(17, 24, 39)
    WHITE = RGBColor(255, 255, 255)

    def add_slide(title, subtitle, script_notes):
        # Using a blank layout to build custom
        slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(slide_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

        # Title Box
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1.0))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = BRAND_GREEN
        p.font.name = "Arial"
        
        # Subtitle Box
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.5))
        stf = sub_box.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(24)
        sp.font.color.rgb = WHITE
        sp.font.name = "Arial"

        # Image Placeholder (Rectangle)
        left = Inches(1.5)
        top = Inches(2.0)
        width = Inches(10.33)
        height = Inches(5.0)
        shape = slide.shapes.add_shape(
            1, left, top, width, height  # 1 = msoShapeRectangle
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(55, 65, 81) # Dark gray placeholder
        shape.line.color.rgb = BRAND_GREEN
        shape.line.width = Pt(2)
        
        # Text inside placeholder
        shape.text = "RIGHT-CLICK HERE -> 'Change Picture' -> 'This Device'\n\n(Insert your screenshot here)"
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].font.size = Pt(24)
        shape.text_frame.paragraphs[0].font.color.rgb = WHITE

        # Speaker Notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = script_notes

        return slide

    # SLIDE 1: Title
    add_slide(
        "SliceMatic: PizzaFlow", 
        "Fast, modern, and data-driven AI ordering system.",
        "Ordering a pizza shouldn't feel like filling out a tax return. It should be fast, intuitive, and delightful. Welcome to PizzaFlow, the custom-built, AI-ready ordering system designed specifically for SliceMatic in New Ashok Nagar. We set out to build a platform that doesn't just take orders—it manages the entire business economics pipeline in real-time."
    )

    # SLIDE 2: Menu
    add_slide(
        "The Customer Experience", 
        "Dynamic Menu & Real-Time Customization",
        "Let's look at the customer experience. We built a completely custom, dynamic interface on top of Gradio. Notice the UI—it’s clean, modern, and perfectly aligned with SliceMatic’s brand identity. Under the hood, this isn't just a static menu. The system dynamically pulls from our backend data structures, meaning if SliceMatic wants to add a new 'Truffle Mushroom' topping tomorrow, the UI updates instantly without touching the frontend code."
    )

    # SLIDE 3: Pricing
    add_slide(
        "The Pricing Engine", 
        "Calculating unit economics, GST, and margins live.",
        "But a pretty interface is only half the battle. The real magic is our custom Pricing Engine. As the user builds their pizza, our backend is instantly calculating unit economics. We've hardcoded the complex Indian tax structure—calculating precise GST, handling promotional discounts, and tracking contribution margins per item. It guarantees that the customer sees exactly what they are paying for, while the business maintains its profit margins."
    )

    # SLIDE 4: Checkout
    add_slide(
        "Checkout & Persistence", 
        "Fault-tolerant storage on Hugging Face.",
        "Checkout is frictionless. But what happens when the customer clicks 'Pay'? This is where PizzaFlow shines. We implemented a robust, fault-tolerant persistence layer. The moment an order is confirmed, it is securely appended to our database and flat-file logs. Because we deployed this using Docker on Hugging Face Spaces with a mounted Persistent Volume, no data is ever lost, even if the server restarts."
    )

    # SLIDE 5: Analytics
    add_slide(
        "The Business Value", 
        "Real-time Analytics Dashboard for the owner.",
        "Finally, we didn't just build an app for the customer; we built a tool for the business owner. Welcome to the Analytics Dashboard. Every single order is instantly parsed and visualized here. The owner can track revenue, filter by payment methods, and download the raw data with a single click. And thanks to our fully automated CI/CD pipeline via GitHub Actions, whenever our developer team pushes a new feature, it is tested and deployed to production with zero downtime."
    )

    # SLIDE 6: Production Readiness
    add_slide(
        "Production Readiness",
        "How we know it's ready to ship.",
        "When asked how we know we're ready for production, we look at it from two perspectives. Engineering: We have automated CI pipelines that block failing tests, zero-touch deployment via GitHub Actions to Hugging Face, environment parity using Docker, and fault-tolerant Persistent Storage. Product: We mapped every feature directly to our PRD, built validation layers for edge cases, verified that our Pricing Engine matches the financial models, and successfully tested the end-to-end user lifecycle. We didn't just build an MVP; we built a production-grade system."
    )

    # SLIDE 7: Conclusion
    add_slide(
        "Thank You", 
        "PizzaFlow by Group 3",
        "PizzaFlow is more than an MVP. It is a highly scalable, data-driven foundation ready for Stage 3's Conversational AI integration. Fast for the customer, profitable for the business. Thank you from Group 3."
    )

    output_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "ppt", "PizzaFlow_Product_Demo.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_demo_ppt()
